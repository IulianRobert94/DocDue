#!/usr/bin/env python3
"""Generate a professional DocTracker presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_BG     = RGBColor(0x0F, 0x0F, 0x14)
DARK_CARD   = RGBColor(0x1C, 0x1C, 0x1E)
GRAY_TEXT   = RGBColor(0xA0, 0xA0, 0xA5)
LIGHT_GRAY  = RGBColor(0x8E, 0x8E, 0x93)
BLUE        = RGBColor(0x00, 0x7A, 0xFF)
GREEN       = RGBColor(0x34, 0xC7, 0x59)
PURPLE      = RGBColor(0xAF, 0x52, 0xDE)
ORANGE      = RGBColor(0xFF, 0x95, 0x00)
RED         = RGBColor(0xFF, 0x3B, 0x30)
ACCENT_BLUE = RGBColor(0x00, 0x6E, 0xE6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, spacing=Pt(8), font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_icon_text(slide, left, top, icon_color, icon_char, title, desc,
                  title_size=20, desc_size=14):
    # Icon circle
    add_circle(slide, left, top + Inches(0.05), Inches(0.5), icon_color)
    add_text_box(slide, left, top + Inches(0.02), Inches(0.5), Inches(0.55),
                 icon_char, font_size=22, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + Inches(0.7), top, Inches(3.5), Inches(0.4),
                 title, font_size=title_size, color=WHITE, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.7), top + Inches(0.35), Inches(3.5),
                 Inches(0.8), desc, font_size=desc_size, color=GRAY_TEXT)

def add_gradient_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    return shape

def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Decorative top accent bar
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE)

# App icon placeholder
icon_shape = add_shape(slide, Inches(5.9), Inches(1.5), Inches(1.5),
                       Inches(1.5), BLUE, corner_radius=0.18)
add_text_box(slide, Inches(5.9), Inches(1.55), Inches(1.5), Inches(1.5),
             "DT", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10.3), Inches(1),
             "DocTracker", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2.5), Inches(4.3), Inches(8.3), Inches(0.6),
             "Smart Document & Expiration Tracking for iOS & Android",
             font_size=22, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Version badge
badge = add_shape(slide, Inches(5.65), Inches(5.2), Inches(2), Inches(0.45),
                  DARK_CARD, corner_radius=0.5)
add_text_box(slide, Inches(5.65), Inches(5.22), Inches(2), Inches(0.45),
             "Version 12.2  •  2025", font_size=13, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_text_box(slide, Inches(3), Inches(6.3), Inches(7.3), Inches(0.5),
             "Never miss a deadline. Track everything that expires.",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
             font_name="Helvetica Neue")

add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), RED)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Problem", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(1.2), Inches(0.05), RED)

# Pain points as cards
pain_points = [
    ("Missed Deadlines", "Expired car insurance, lapsed IDs, overdue bills — people forget critical document dates and face fines, penalties, or gaps in coverage.", RED),
    ("Scattered Records", "Documents tracked across paper folders, phone notes, email reminders, and mental notes. No single source of truth.", ORANGE),
    ("No Smart Alerts", "Calendar reminders are manual and static. No automatic awareness of which documents are approaching expiration.", PURPLE),
    ("Financial Surprises", "No visibility into upcoming payment obligations across insurance, utilities, subscriptions, and contracts.", BLUE),
]

for i, (title, desc, color) in enumerate(pain_points):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.1)
    y = Inches(1.8) + row * Inches(2.5)

    card = add_shape(slide, x, y, Inches(5.7), Inches(2.1), DARK_CARD, 0.06)

    # Accent bar on the left of card
    add_gradient_bar(slide, x + Inches(0.15), y + Inches(0.3),
                     Inches(0.06), Inches(1.5), color)

    add_text_box(slide, x + Inches(0.5), y + Inches(0.25), Inches(4.8),
                 Inches(0.4), title, font_size=22, color=color, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.75), Inches(4.8),
                 Inches(1.2), desc, font_size=15, color=GRAY_TEXT)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Solution — DocTracker", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(2.2), Inches(0.05), GREEN)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
             "One app to track every document, contract, bill, and subscription that expires.",
             font_size=20, color=GRAY_TEXT)

features = [
    ("Centralized Tracking", "All documents in one place — vehicle, personal, home, financial. Organized by category with smart subtypes.", BLUE, "📋"),
    ("Intelligent Alerts", "3-tier status system (OK → Due Soon → Expired) with customizable warning thresholds per category.", GREEN, "🔔"),
    ("Automatic Reminders", "Push notifications at 9 AM, days before expiration. Configurable: 1, 3, 7, or 14 days ahead.", ORANGE, "⏰"),
    ("Offline & Private", "100% local storage. Zero cloud dependency. Your data never leaves your device.", PURPLE, "🔒"),
]

for i, (title, desc, color, icon) in enumerate(features):
    y = Inches(2.5) + i * Inches(1.2)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.0), DARK_CARD, 0.06)
    add_circle(slide, Inches(1.1), y + Inches(0.2), Inches(0.6), color)
    add_text_box(slide, Inches(1.1), y + Inches(0.15), Inches(0.6), Inches(0.6),
                 icon, font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), y + Inches(0.12), Inches(4), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.52), Inches(10), Inches(0.5),
                 desc, font_size=14, color=GRAY_TEXT)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Four Categories
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Four Smart Categories", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.05), BLUE)

categories = [
    ("Vehicles", BLUE, "🚗",
     ["RCA (Car Insurance)", "ITP (Inspection)", "CASCO", "Road Tax / Vignette",
      "Service Records", "Tahograph", "ADR Verification"]),
    ("Home", GREEN, "🏠",
     ["Electricity, Gas, Water", "Internet & Phone", "Rent Contracts",
      "Home Insurance (PAD)", "Subscriptions (Streaming, Cloud)", "Property Tax"]),
    ("Personal", PURPLE, "👤",
     ["National ID & Passport", "Driver's License", "Medical Records",
      "Professional Certificates", "Employment Contracts", "Life & Health Insurance"]),
    ("Financial", ORANGE, "💼",
     ["Business Licenses", "Client Contracts", "Tax Declarations (ANAF)",
      "Credit Payments", "Leasing", "Fines & Penalties"]),
]

for i, (name, color, icon, subtypes) in enumerate(categories):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(1.7)

    card = add_shape(slide, x, y, Inches(2.95), Inches(5.3), DARK_CARD, 0.05)

    # Colored top accent
    add_gradient_bar(slide, x + Inches(0.15), y + Inches(0.15),
                     Inches(2.65), Inches(0.06), color)

    # Icon + Title
    add_text_box(slide, x, y + Inches(0.4), Inches(2.95), Inches(0.5),
                 f"{icon}  {name}", font_size=24, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtypes
    for j, sub in enumerate(subtypes):
        add_text_box(slide, x + Inches(0.3), y + Inches(1.1) + j * Inches(0.52),
                     Inches(2.5), Inches(0.45),
                     f"•  {sub}", font_size=13, color=GRAY_TEXT)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Core Features Deep Dive
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Core Features", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(1.6), Inches(0.05), PURPLE)

features_grid = [
    ("OCR Document Scanning", "Point your camera at a document → ML Kit extracts dates, amounts, and document type automatically. Fill forms in seconds.", "📸", BLUE),
    ("Smart Notifications", "Local push notifications at 9:00 AM, scheduled 1/3/7/14 days before expiration. Snooze or mark as renewed directly from the notification.", "🔔", GREEN),
    ("Home Screen Widget", "At-a-glance view of expired and upcoming documents without opening the app. Shows next 3 deadlines with countdown.", "📱", PURPLE),
    ("Calendar Integration", "One-tap to add any document's due date to your device calendar. Works with Apple Calendar and Google Calendar.", "📅", ORANGE),
    ("Recurrence Engine", "Auto-advance due dates when marked as paid. Supports: monthly, quarterly, biannual, annual, and custom recurrence.", "🔄", BLUE),
    ("Biometric Lock", "Secure your data with Face ID or fingerprint authentication. Optional app-level security for sensitive documents.", "🔐", RED),
]

for i, (title, desc, icon, color) in enumerate(features_grid):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.15)
    y = Inches(1.7) + row * Inches(2.7)

    card = add_shape(slide, x, y, Inches(3.9), Inches(2.35), DARK_CARD, 0.05)

    add_circle(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.55), color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.22), Inches(0.55),
                 Inches(0.55), icon, font_size=22, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.0), y + Inches(0.28), Inches(2.7),
                 Inches(0.35), title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.95), Inches(3.4),
                 Inches(1.2), desc, font_size=13, color=GRAY_TEXT)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — User Flow
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "User Journey", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), ORANGE)

steps = [
    ("1", "Onboard", "Choose categories\nyou want to track", BLUE),
    ("2", "Add Docs", "Quick form or\nOCR scan photo", GREEN),
    ("3", "Get Alerts", "Smart notifications\nbefore expiry", ORANGE),
    ("4", "Take Action", "Renew, pay, or\nextend documents", PURPLE),
    ("5", "Stay Safe", "Never miss a\ndeadline again", RED),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(2.5)
    y = Inches(2.2)

    # Connector arrow (except last)
    if i < len(steps) - 1:
        arrow_x = x + Inches(2.2)
        add_text_box(slide, arrow_x, y + Inches(0.6), Inches(0.5), Inches(0.4),
                     "→", font_size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    card = add_shape(slide, x, y, Inches(2.2), Inches(2.8), DARK_CARD, 0.06)

    # Step number circle
    add_circle(slide, x + Inches(0.75), y + Inches(0.3), Inches(0.7), color)
    add_text_box(slide, x + Inches(0.75), y + Inches(0.3), Inches(0.7),
                 Inches(0.7), num, font_size=28, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x, y + Inches(1.2), Inches(2.2), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.65), Inches(1.9),
                 Inches(0.9), desc, font_size=14, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)

# Bottom feature strip
strip_y = Inches(5.5)
strip_items = [
    "5-Tab Navigation", "Pull-to-Refresh", "Debounced Search",
    "Multi-Filter System", "Animated Transitions"
]
for i, item in enumerate(strip_items):
    x = Inches(0.5) + i * Inches(2.5)
    pill = add_shape(slide, x, strip_y, Inches(2.2), Inches(0.45), DARK_CARD, 0.5)
    add_text_box(slide, x, strip_y + Inches(0.03), Inches(2.2), Inches(0.4),
                 item, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Alert System
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Intelligent Alert System", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(2.2), Inches(0.05), GREEN)

# Three status cards
statuses = [
    ("EXPIRED", "Past due date", "Documents that have already expired and need immediate attention.", RED, "✕"),
    ("DUE SOON", "Within warning threshold", "Approaching expiration — configurable per category (default: 14 days).", ORANGE, "!"),
    ("OK", "Not yet due", "Documents that are current and not approaching their expiration date.", GREEN, "✓"),
]

for i, (status, subtitle, desc, color, icon) in enumerate(statuses):
    x = Inches(0.6) + i * Inches(4.15)
    y = Inches(1.8)
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.0), DARK_CARD, 0.05)
    add_gradient_bar(slide, x, y, Inches(3.9), Inches(0.08), color)

    add_circle(slide, x + Inches(0.3), y + Inches(0.35), Inches(0.55), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.32), Inches(0.55),
                 Inches(0.55), icon, font_size=24, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.05), y + Inches(0.35), Inches(2.5),
                 Inches(0.35), status, font_size=22, color=color, bold=True)
    add_text_box(slide, x + Inches(1.05), y + Inches(0.7), Inches(2.5),
                 Inches(0.3), subtitle, font_size=13, color=LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.15), Inches(3.3),
                 Inches(0.7), desc, font_size=13, color=GRAY_TEXT)

# Notification details section
y2 = Inches(4.2)
add_text_box(slide, Inches(0.8), y2, Inches(11), Inches(0.5),
             "Notification Engine", font_size=24, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), y2 + Inches(0.45), Inches(1.6), Inches(0.04), GREEN)

notif_features = [
    ("Scheduled at 9:00 AM", "Respectful timing — alerts arrive during morning routine, not at midnight."),
    ("Configurable Lead Time", "Choose reminder days: 1, 3, 7, or 14 days before expiration."),
    ("Quick Actions", "\"Renewed ✓\" or \"Tomorrow\" snooze — respond directly from the notification."),
    ("Per-Category Thresholds", "RCA needs 7 days notice, rent needs 14 days — fully customizable per category."),
]

for i, (title, desc) in enumerate(notif_features):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.1)
    y = y2 + Inches(0.7) + row * Inches(1.1)
    add_text_box(slide, x, y, Inches(5.5), Inches(0.3),
                 f"→  {title}", font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.35), y + Inches(0.35), Inches(5.3),
                 Inches(0.5), desc, font_size=13, color=GRAY_TEXT)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Technology Stack", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.05), ACCENT_BLUE)

tech_cols = [
    ("Framework", BLUE, [
        ("React Native 0.81", "Cross-platform native UI"),
        ("Expo SDK 54", "Managed workflow + dev tools"),
        ("Expo Router 6", "File-based navigation"),
        ("TypeScript 5.9", "Full type safety"),
    ]),
    ("State & Data", GREEN, [
        ("Zustand 5.0", "Lightweight state management"),
        ("AsyncStorage", "Offline-first persistence"),
        ("JSON / Excel Backup", "Import & export support"),
        ("Auto-Migration", "Seamless version upgrades"),
    ]),
    ("Services", ORANGE, [
        ("ML Kit OCR", "On-device text recognition"),
        ("Local Notifications", "Scheduled push alerts"),
        ("Expo Calendar", "Device calendar integration"),
        ("Expo Widgets", "Home screen widgets"),
    ]),
    ("UX & Security", PURPLE, [
        ("Reanimated 4.1", "60fps native animations"),
        ("Haptic Feedback", "Tactile UI interactions"),
        ("Face ID / Touch ID", "Biometric app lock"),
        ("Apple HIG", "iOS design compliance"),
    ]),
]

for i, (col_title, color, items) in enumerate(tech_cols):
    x = Inches(0.5) + i * Inches(3.2)
    y = Inches(1.7)

    card = add_shape(slide, x, y, Inches(3.0), Inches(5.0), DARK_CARD, 0.05)
    add_gradient_bar(slide, x, y, Inches(3.0), Inches(0.07), color)

    add_text_box(slide, x, y + Inches(0.25), Inches(3.0), Inches(0.4),
                 col_title, font_size=20, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    for j, (tech, desc) in enumerate(items):
        ty = y + Inches(0.85) + j * Inches(1.0)
        add_text_box(slide, x + Inches(0.25), ty, Inches(2.5), Inches(0.3),
                     tech, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.25), ty + Inches(0.3), Inches(2.5),
                     Inches(0.4), desc, font_size=12, color=GRAY_TEXT)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Architecture & Data Flow", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(2.4), Inches(0.05), PURPLE)

# Data flow pipeline
flow_steps = [
    ("RawDocument", "AsyncStorage\n(dt12_docs)", BLUE),
    ("enrichDocument()", "Add _daysUntil\n& _status", GREEN),
    ("EnrichedDocument", "Runtime computed\nfields", ORANGE),
    ("UI Screens", "Consumed via\nhooks & selectors", PURPLE),
]

for i, (title, desc, color) in enumerate(flow_steps):
    x = Inches(0.5) + i * Inches(3.2)
    y = Inches(1.8)

    card = add_shape(slide, x, y, Inches(2.8), Inches(1.6), DARK_CARD, 0.06)
    add_gradient_bar(slide, x, y, Inches(2.8), Inches(0.06), color)

    add_text_box(slide, x, y + Inches(0.2), Inches(2.8), Inches(0.4),
                 title, font_size=18, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.65), Inches(2.8), Inches(0.7),
                 desc, font_size=13, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    if i < len(flow_steps) - 1:
        add_text_box(slide, x + Inches(2.85), y + Inches(0.55), Inches(0.4),
                     Inches(0.4), "→", font_size=24, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

# Architecture highlights
y3 = Inches(3.8)
add_text_box(slide, Inches(0.8), y3, Inches(11), Inches(0.4),
             "Key Architecture Decisions", font_size=22, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), y3 + Inches(0.4), Inches(2), Inches(0.04), PURPLE)

arch_items = [
    ("Offline-First Design", "All data persisted locally via AsyncStorage. No server needed. Works without internet."),
    ("Ephemeral Computed Fields", "Status and countdown are calculated at runtime with _ prefix. Stripped before saving to prevent stale data."),
    ("Auto-Persisting Stores", "Every Zustand mutation automatically saves to AsyncStorage and reschedules notifications."),
    ("Safe Date Handling", "Custom parseLocalDate() prevents UTC timezone bugs. Never uses new Date(\"YYYY-MM-DD\")."),
    ("Two-Store Pattern", "Documents store (CRUD + enrichment) and Settings store (preferences + theme) — separated concerns."),
    ("Auto-Migration", "New settings fields merge automatically via spread pattern. No manual migration scripts needed."),
]

for i, (title, desc) in enumerate(arch_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.2)
    y = y3 + Inches(0.7) + row * Inches(1.05)
    add_text_box(slide, x, y, Inches(5.8), Inches(0.3),
                 f"▸  {title}", font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(5.5),
                 Inches(0.5), desc, font_size=12, color=GRAY_TEXT)

add_slide_number(slide, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Bilingual & UX
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Bilingual UX & Design System", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(2.8), Inches(0.05), BLUE)

# Left column: i18n
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
             "Romanian & English", font_size=24, color=BLUE, bold=True)

i18n_items = [
    "300+ translation keys with placeholder interpolation",
    "Subtype-aware: \"RCA\" stays as-is; English shows \"Car Insurance\"",
    "Fallback chain: requested language → Romanian → raw key",
    "All UI strings: navigation, forms, errors, confirmations, notifications",
    "Custom i18n engine — no external library dependency",
    "One-tap language switch in Settings",
]
for i, item in enumerate(i18n_items):
    add_text_box(slide, Inches(1.0), Inches(2.3) + i * Inches(0.55),
                 Inches(5.2), Inches(0.45),
                 f"•  {item}", font_size=14, color=GRAY_TEXT)

# Right column: Design
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
             "Apple HIG Design System", font_size=24, color=PURPLE, bold=True)

design_items = [
    ("Dark & Light Themes", "Pure black dark mode, grouped light mode"),
    ("Large Title Headers", "Scroll-driven collapse into compact bar"),
    ("Semantic Colors", "Apple system colors + status indicators"),
    ("Native Controls", "Segmented controls, grouped tables, sheets"),
    ("Spring Animations", "Scale 0.97 on press + bounce feedback"),
    ("Haptic Feedback", "Light, medium, selection — tactile response"),
]
for i, (title, desc) in enumerate(design_items):
    y = Inches(2.3) + i * Inches(0.72)
    add_text_box(slide, Inches(7.2), y, Inches(5.5), Inches(0.3),
                 f"▸  {title}", font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.5), y + Inches(0.28), Inches(5.3),
                 Inches(0.35), desc, font_size=12, color=GRAY_TEXT)

# Theme cards at bottom
theme_cards = [
    ("Dark Mode", "#000000 bg, #1C1C1E cards", DARK_CARD),
    ("Status Colors", "Red / Orange / Green", RGBColor(0x2C, 0x2C, 0x2E)),
    ("Accent Blue", "#007AFF — primary action", RGBColor(0x2C, 0x2C, 0x2E)),
]
for i, (name, desc, bg) in enumerate(theme_cards):
    x = Inches(0.8) + i * Inches(4.15)
    y = Inches(6.2)
    card = add_shape(slide, x, y, Inches(3.9), Inches(0.8), bg, 0.08)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.08), Inches(3.3),
                 Inches(0.35), name, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.4), Inches(3.3),
                 Inches(0.35), desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Monetization & Premium
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Business Model & Roadmap", font_size=40, color=WHITE, bold=True)
add_gradient_bar(slide, Inches(0.8), Inches(1.2), Inches(2.4), Inches(0.05), ORANGE)

# Free vs Premium comparison
for i, (tier, color, features) in enumerate([
    ("Free Tier", LIGHT_GRAY, [
        "Up to 5 documents",
        "All 4 categories",
        "Basic status tracking",
        "Search & filter",
        "JSON backup/restore",
    ]),
    ("Premium", ORANGE, [
        "Unlimited documents",
        "Push notifications",
        "OCR document scanning",
        "Spending analytics",
        "Excel import/export",
        "Home screen widgets",
        "Priority support",
    ]),
]):
    x = Inches(0.8) + i * Inches(4.3)
    y = Inches(1.7)
    w = Inches(4.0)
    h = Inches(4.5) if i == 1 else Inches(3.8)

    card = add_shape(slide, x, y, w, h, DARK_CARD, 0.05)
    add_gradient_bar(slide, x, y, w, Inches(0.07), color)

    add_text_box(slide, x, y + Inches(0.25), w, Inches(0.4),
                 tier, font_size=24, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    for j, feat in enumerate(features):
        add_text_box(slide, x + Inches(0.4), y + Inches(0.85) + j * Inches(0.42),
                     Inches(3.2), Inches(0.35),
                     f"✓  {feat}", font_size=14,
                     color=WHITE if i == 1 else GRAY_TEXT)

# Roadmap on the right
rx = Inches(9.3)
ry = Inches(1.7)
add_text_box(slide, rx, ry, Inches(3.5), Inches(0.4),
             "Roadmap", font_size=24, color=BLUE, bold=True)
add_gradient_bar(slide, rx, ry + Inches(0.4), Inches(1), Inches(0.04), BLUE)

roadmap = [
    ("Q1 2025", "Core app launch\niOS + Android", GREEN),
    ("Q2 2025", "Premium tier\n+ OCR scanning", ORANGE),
    ("Q3 2025", "Fleet management\n+ multi-user", BLUE),
    ("Q4 2025", "Cloud sync\n+ web dashboard", PURPLE),
]

for i, (quarter, desc, color) in enumerate(roadmap):
    y = ry + Inches(0.65) + i * Inches(1.25)
    add_circle(slide, rx, y + Inches(0.05), Inches(0.35), color)
    # Connector line (vertical)
    if i < len(roadmap) - 1:
        add_gradient_bar(slide, rx + Inches(0.14), y + Inches(0.42),
                         Inches(0.06), Inches(0.9), DARK_CARD)
    add_text_box(slide, rx + Inches(0.5), y, Inches(3), Inches(0.3),
                 quarter, font_size=15, color=color, bold=True)
    add_text_box(slide, rx + Inches(0.5), y + Inches(0.3), Inches(3),
                 Inches(0.6), desc, font_size=12, color=GRAY_TEXT)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You / CTA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE)

# App icon
icon_shape = add_shape(slide, Inches(5.9), Inches(1.3), Inches(1.5),
                       Inches(1.5), BLUE, corner_radius=0.18)
add_text_box(slide, Inches(5.9), Inches(1.35), Inches(1.5), Inches(1.5),
             "DT", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.8),
             "Thank You", font_size=52, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.5), Inches(4.0), Inches(8.3), Inches(0.6),
             "DocTracker — Smart Document & Expiration Tracking",
             font_size=22, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Key stats strip
stats = [
    ("4", "Categories"),
    ("50+", "Document Types"),
    ("2", "Languages"),
    ("100%", "Offline"),
    ("0", "Cloud Required"),
]

for i, (number, label) in enumerate(stats):
    x = Inches(1.1) + i * Inches(2.3)
    y = Inches(5.0)
    card = add_shape(slide, x, y, Inches(2.0), Inches(1.2), DARK_CARD, 0.08)
    add_text_box(slide, x, y + Inches(0.15), Inches(2.0), Inches(0.5),
                 number, font_size=32, color=BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.65), Inches(2.0), Inches(0.4),
                 label, font_size=13, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(6.5), Inches(7.3), Inches(0.5),
             "Questions & Discussion", font_size=18, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 12, TOTAL_SLIDES)


# ── Save ────────────────────────────────────────────────────────────────
output_path = "/Users/andreiiulian/Desktop/doctracker V12.2 Code test/DocTracker_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
